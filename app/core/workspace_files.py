"""Workspace file management — index, upload, read, search, Office parsing.

Independent of DataStore backend: file index is always JSON on disk.
FileManager operates on the workspace's files/ directory under data_dir.
"""

from __future__ import annotations

import json
import fnmatch
from pathlib import Path
from typing import Any

from loguru import logger

from app.config import config
from app.core.file_handler import (
    MAX_FILE_SIZE,
    ALLOWED_EXTENSIONS,
    FileValidationError,
    decode_content,
    validate_upload,
)

# ── Office parsing (optional) ──

_HAS_DOCX = False
try:
    import docx as _docx_module
    _HAS_DOCX = True
except ImportError:
    pass


def parse_docx(file_path: str | Path) -> str:
    if not _HAS_DOCX:
        raise RuntimeError("python-docx not installed. Run: pip install python-docx")
    doc = _docx_module.Document(str(file_path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def parse_xlsx(file_path: str | Path) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(file_path), read_only=True)
    lines = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"=== {sheet_name} ===")
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(c) if c is not None else "" for c in row)
            if row_text.strip(" |"):
                lines.append(row_text)
    wb.close()
    return "\n".join(lines)


# ── PPTX parsing (optional) ──

_HAS_PPTX = False
try:
    import pptx as _pptx_module
    _HAS_PPTX = True
except ImportError:
    pass


def parse_pptx(file_path: str | Path) -> str:
    """Extract all text from .pptx slides."""
    if not _HAS_PPTX:
        raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")
    prs = _pptx_module.Presentation(str(file_path))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text for cell in row.cells]
                    lines.append(" | ".join(cells))
    return "\n".join(lines)


_HAS_PIL = False
try:
    from PIL import Image as _PIL_Image
    _HAS_PIL = True
except ImportError:
    pass


def parse_image(file_path: str | Path) -> str:
    """Extract metadata from an image file."""
    if not _HAS_PIL:
        raise RuntimeError("Pillow not installed. Run: pip install Pillow")
    img = _PIL_Image.open(str(file_path))
    return (
        f"[图片: {Path(file_path).name}]\n"
        f"格式: {img.format}\n"
        f"尺寸: {img.width}x{img.height}\n"
        f"模式: {img.mode}\n"
        f"\n如需分析图片内容，请使用 LLM 的视觉分析能力。"
    )


def is_text_file(ext: str) -> bool:
    return ext in ("md", "txt", "json", "yaml", "yml")


def is_image_ext(ext: str) -> bool:
    return ext in ("png", "jpg", "jpeg", "gif", "bmp")


def _validate_file_path(base_dir: Path, relative_path: str) -> Path:
    resolved = (base_dir / relative_path).resolve()
    if not str(resolved).startswith(str(base_dir.resolve())):
        raise FileValidationError("Path traversal detected")
    return resolved


def _index_path(files_dir: Path) -> Path:
    return files_dir / "_index.json"


def _load_index(files_dir: Path) -> list[dict]:
    p = _index_path(files_dir)
    if not p.exists():
        return []
    try:
        return json.loads(p.read_text(encoding="utf-8"))
    except (json.JSONDecodeError, OSError):
        logger.warning(f"Corrupted or unreadable index file: {p}")
        return []


def _save_index(files_dir: Path, index: list[dict]) -> None:
    p = _index_path(files_dir)
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(index, ensure_ascii=False, indent=2), encoding="utf-8")


def _now() -> str:
    from datetime import datetime, timezone
    return datetime.now(timezone.utc).isoformat()


class WorkspaceFileManager:
    """Manages files within a single workspace. File path: <data_dir>/workspaces/<ws_id>/files/"""

    def __init__(self, workspace_id: str):
        ws_dir = Path(config.data_dir) / "workspaces" / workspace_id
        self._files_dir = ws_dir / "files"
        self._uploads_dir = self._files_dir / "uploads"
        self._generated_dir = self._files_dir / "_generated"

    # ── public API ──

    def list_files(self, pattern: str = "*", max_results: int = 50) -> list[dict]:
        index = _load_index(self._files_dir)
        if pattern != "*":
            matched = [f for f in index if fnmatch.fnmatch(f["path"], pattern)]
        else:
            matched = list(index)
        matched.sort(key=lambda f: f.get("uploaded_at", ""), reverse=True)
        return matched[:max_results]

    def get_file_info(self, relative_path: str) -> dict | None:
        safe = Path(relative_path).name
        for f in _load_index(self._files_dir):
            if f["path"] == safe:
                return f
        return None

    def read_file(self, relative_path: str, max_chars: int = 8000) -> dict:
        entry = self.get_file_info(relative_path)
        if entry is None:
            raise FileNotFoundError(f"File not found: {relative_path}")

        ext = entry.get("type", "")
        if entry["source"] in ("upload", "generated"):
            base = self._uploads_dir if entry["source"] == "upload" else self._generated_dir
            file_path = _validate_file_path(base, relative_path)
        elif entry.get("abs_path"):
            file_path = Path(entry["abs_path"])
            if not file_path.exists():
                raise FileNotFoundError(f"Linked file not found: {file_path}")
            if not str(file_path.resolve()).startswith(str(Path(config.data_dir).resolve())):
                raise FileValidationError("Linked file path outside allowed directory")
        else:
            raise FileNotFoundError(f"Cannot resolve file: {relative_path}")

        size = file_path.stat().st_size
        if is_text_file(ext):
            text = file_path.read_text(encoding="utf-8", errors="replace")
        elif ext == "docx":
            text = parse_docx(file_path)
        elif ext == "xlsx":
            try:
                text = parse_xlsx(file_path)
            except Exception:
                text = f"[无法解析 xlsx 文件: {relative_path}]"
        elif ext == "pptx":
            try:
                text = parse_pptx(file_path)
            except Exception:
                text = f"[无法解析 pptx 文件: {relative_path}]"
        elif is_image_ext(ext):
            try:
                text = parse_image(file_path)
            except Exception:
                text = f"[无法读取图片: {relative_path}]"
        else:
            text = file_path.read_text(encoding="utf-8", errors="replace")

        truncated = False
        if len(text) > max_chars:
            text = text[:max_chars]
            truncated = True
        return {"path": relative_path, "text": text, "size": size, "type": ext, "truncated": truncated}

    def read_file_section(self, relative_path: str, start: int = 1, end: int = 100) -> dict:
        content = self.read_file(relative_path, max_chars=10 * 1024 * 1024)
        lines = content["text"].split("\n")
        return {
            "path": relative_path,
            "lines": lines[start - 1 : end],
            "start_line": start, "end_line": min(end, len(lines)),
            "total_lines": len(lines),
        }

    def search_files(self, query: str, file_pattern: str = "*.md", max_results: int = 10) -> list[dict]:
        results = []
        for entry in _load_index(self._files_dir):
            if not fnmatch.fnmatch(entry["path"], file_pattern):
                continue
            try:
                content = self.read_file(entry["path"], max_chars=100 * 1024)
            except (FileNotFoundError, RuntimeError):
                continue
            lines = content["text"].split("\n")
            matches = [(i + 1, l.strip()[:200]) for i, l in enumerate(lines) if query.lower() in l.lower()]
            if matches:
                results.append({
                    "path": entry["path"],
                    "type": entry.get("type", ""),
                    "matches": [{"line": ln, "text": txt} for ln, txt in matches[:10]],
                    "match_count": len(matches),
                })
            if len(results) >= max_results:
                break
        return results

    def write_file(self, relative_path: str, content: str) -> dict:
        if not relative_path or not relative_path.strip():
            raise ValueError("File path cannot be empty")
        safe_name = Path(relative_path).name
        dest = self._generated_dir / safe_name
        dest.parent.mkdir(parents=True, exist_ok=True)
        dest.write_text(content, encoding="utf-8")

        self._upsert_index({
            "path": safe_name,
            "size": len(content.encode("utf-8")),
            "type": Path(safe_name).suffix.lower().lstrip("."),
            "source": "generated",
            "uploaded_at": _now(),
            "uploaded_by": "agent",
            "summary": content[:100].replace("\n", " "),
        })
        logger.info(f"[ws files] Generated: {safe_name}")
        return {"path": safe_name, "size": len(content)}

    def upload_file(self, filename: str, content: bytes, uploaded_by: str = "") -> dict:
        validate_upload(filename, content)
        safe_name = Path(filename).name
        ext = safe_name.rsplit(".", 1)[-1].lower() if "." in safe_name else "bin"
        dest = self._uploads_dir / safe_name
        dest.parent.mkdir(parents=True, exist_ok=True)
        if not (dest.exists() and dest.stat().st_size == len(content)):
            dest.write_bytes(content)

        # Parse for summary
        text = decode_content(content) if is_text_file(ext) else ""
        if ext == "docx" and _HAS_DOCX:
            try:
                text = parse_docx(dest)
            except RuntimeError:
                text = ""
        elif ext == "xlsx":
            try:
                text = parse_xlsx(dest)
            except Exception:
                text = ""
        summary = text[:100].replace("\n", " ") if text else ""
        if not summary:
            summary = f"Uploaded {filename}"

        self._upsert_index({
            "path": safe_name,
            "size": len(content),
            "type": ext,
            "source": "upload",
            "uploaded_at": _now(),
            "uploaded_by": uploaded_by,
            "summary": summary,
        })
        return {"path": safe_name, "size": len(content), "type": ext}

    def delete_file(self, relative_path: str) -> bool:
        safe_name = Path(relative_path).name
        index = _load_index(self._files_dir)
        entries = [f for f in index if f["path"] == safe_name]
        if not entries:
            return False
        index = [f for f in index if f["path"] != safe_name]
        _save_index(self._files_dir, index)

        entry = entries[0]
        if entry["source"] == "upload":
            target = self._uploads_dir / safe_name
        elif entry["source"] == "generated":
            target = self._generated_dir / safe_name
        else:
            target = None
        if target and target.exists():
            target.unlink()
        return True

    def get_info(self) -> dict:
        index = _load_index(self._files_dir)
        types, sources = {}, {}
        for f in index:
            types[f.get("type", "?")] = types.get(f.get("type", "?"), 0) + 1
            sources[f.get("source", "?")] = sources.get(f.get("source", "?"), 0) + 1
        return {"file_count": len(index), "by_type": types, "by_source": sources,
                "files": [{"path": f["path"], "type": f.get("type", ""), "size": f.get("size", 0),
                           "source": f.get("source", ""), "uploaded_at": f.get("uploaded_at", ""),
                           "summary": f.get("summary", "")} for f in index]}

    def _upsert_index(self, entry: dict) -> None:
        index = _load_index(self._files_dir)
        index = [f for f in index if f["path"] != entry["path"]]
        index.append(entry)
        _save_index(self._files_dir, index)
